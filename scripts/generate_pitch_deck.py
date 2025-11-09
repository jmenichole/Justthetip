#!/usr/bin/env python3
"""
JustTheTip - Pitch Deck Generator

Generate a 3-slide pitch deck from a Markdown summary.

Copyright (c) 2025 JustTheTip Bot

This file is part of JustTheTip.

Licensed under the JustTheTip Custom License (Based on MIT).
See LICENSE file in the project root for full license information.

SPDX-License-Identifier: MIT

This software may not be sold commercially without permission.
"""

import argparse
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

DEFAULT_TITLES = ["Problem", "Solution", "Architecture"]


def load_markdown_lines(markdown_path: Path) -> list[str]:
    text = markdown_path.read_text(encoding='utf-8')
    document = Document()
    for line in text.splitlines():
        document.add_paragraph(line.strip())
    return [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]


def bucket_content(paragraphs: list[str]) -> dict[str, list[str]]:
    sections: dict[str, list[str]] = {title: [] for title in DEFAULT_TITLES}
    current = DEFAULT_TITLES[0]

    for entry in paragraphs:
        normalized = entry.strip()
        if normalized.startswith('#'):
            heading = normalized.lstrip('#').strip().title()
            if heading in sections:
                current = heading
            continue

        sections[current].append(normalized)

    return sections


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for bullet in bullets or ["Add supporting details here."]:
        p = body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.level = 0


def build_deck(markdown_path: Path, output_path: Path) -> None:
    paragraphs = load_markdown_lines(markdown_path)
    sections = bucket_content(paragraphs)

    presentation = Presentation()
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)

    for title in DEFAULT_TITLES:
        add_slide(presentation, title, sections.get(title, []))

    presentation.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description='Generate a pitch deck from Markdown content.')
    parser.add_argument('markdown', type=Path, help='Path to the Markdown summary file')
    parser.add_argument('--output', type=Path, default=Path('justthetip_pitch.pptx'), help='Output PowerPoint file')
    args = parser.parse_args()

    build_deck(args.markdown, args.output)
    print(f"✅ Pitch deck created at {args.output}")


if __name__ == '__main__':
    main()
